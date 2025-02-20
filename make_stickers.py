import argparse
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import os


def create_sticker_grid(x, y, folder_name, output_file):
    # Path to the folder containing the images
    folder_path = os.path.join("raw_pictures", folder_name)

    # Check if the folder exists
    if not os.path.exists(folder_path):
        print(f"Folder '{folder_path}' does not exist.")
        return

    # Get all image files from the folder
    image_files = [
        os.path.join(folder_path, f)
        for f in os.listdir(folder_path)
        if os.path.isfile(os.path.join(folder_path, f))
        and f.lower().endswith((".png", ".jpg", ".jpeg", ".gif", ".bmp"))
    ]

    if not image_files:
        print(f"No images found in '{folder_path}'.")
        return

    # Create a new presentation
    prs = Presentation()
    prs.slide_width = Inches(11.69)
    prs.slide_height = Inches(8.27)

    # Calculate the width and height of each square box
    slide_width = Inches(11.69)
    slide_height = Inches(8.27)
    box_width = slide_width / x
    box_height = slide_height / y

    # Initialize variables for tracking images and slides
    total_images = len(image_files)
    current_image_index = 0

    while current_image_index < total_images:
        # Add a new slide
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout

        # Add images to the slide
        for i in range(y):
            for j in range(x):
                if current_image_index >= total_images:
                    break  # Stop if we run out of images

                img_path = image_files[current_image_index]
                try:
                    img = Image.open(img_path)
                except Exception as e:
                    print(f"Failed to open image {img_path}: {e}")
                    current_image_index += 1
                    continue

                # Calculate the position
                left = j * box_width
                top = i * box_height

                # Calculate the aspect ratio
                img_ratio = img.width / img.height
                box_ratio = box_width / box_height

                if img_ratio > box_ratio:
                    # Image is wider than the box
                    new_width = box_width
                    new_height = box_width / img_ratio
                else:
                    # Image is taller than the box
                    new_height = box_height
                    new_width = box_height * img_ratio

                # Calculate the position to center the image
                left_offset = (box_width - new_width) / 2
                top_offset = (box_height - new_height) / 2

                # Add the image to the slide
                slide.shapes.add_picture(
                    img_path,
                    left + left_offset,
                    top + top_offset,
                    width=new_width,
                    height=new_height,
                )

                # Move to the next image
                current_image_index += 1

    # Save the presentation
    prs.save(output_file)
    print(f"Presentation saved as {output_file}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Create a grid of stickers in a .pptx file using images from a folder."
    )
    parser.add_argument("x", type=int, help="Number of columns in the grid")
    parser.add_argument("y", type=int, help="Number of rows in the grid")
    parser.add_argument(
        "folder_name",
        type=str,
        help="Name of the folder inside 'raw_pictures' containing the images",
    )
    parser.add_argument(
        "output_file", type=str, help="Output file name (e.g., output.pptx)"
    )

    args = parser.parse_args()

    create_sticker_grid(args.x, args.y, args.folder_name, args.output_file)
